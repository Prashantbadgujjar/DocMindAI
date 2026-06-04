
from dataclasses import dataclass
from io import BytesIO

from docx import Document
from pptx import Presentation
from pypdf import PdfReader


@dataclass
class TextBlock:
    text: str
    source: str
    location: str


def load_documents(files):

    blocks = []

    for file in files:

        name = file.name.lower()

        if name.endswith(".pdf"):
            blocks.extend(read_pdf(file))

        elif name.endswith(".pptx"):
            blocks.extend(read_ppt(file))

        elif name.endswith(".docx"):
            blocks.extend(read_docx(file))

        elif name.endswith(".txt"):
            blocks.extend(read_txt(file))

    return blocks


def read_pdf(file):

    reader = PdfReader(file)

    blocks = []

    for i, page in enumerate(reader.pages):

        text = page.extract_text()

        if text:
            blocks.append(
                TextBlock(
                    text=text,
                    source=file.name,
                    location=f"Page {i+1}"
                )
            )

    return blocks


def read_ppt(file):

    prs = Presentation(BytesIO(file.read()))

    blocks = []

    for i, slide in enumerate(prs.slides):

        slide_text = []

        for shape in slide.shapes:

            if hasattr(shape, "text"):
                slide_text.append(shape.text)

        text = "\n".join(slide_text)

        if text.strip():
            blocks.append(
                TextBlock(
                    text=text,
                    source=file.name,
                    location=f"Slide {i+1}"
                )
            )

    return blocks


def read_docx(file):

    doc = Document(file)

    text = "\n".join([p.text for p in doc.paragraphs])

    return [
        TextBlock(
            text=text,
            source=file.name,
            location="Document"
        )
    ]


def read_txt(file):

    text = file.read().decode("utf-8")

    return [
        TextBlock(
            text=text,
            source=file.name,
            location="Text File"
        )
    ]
